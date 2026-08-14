from __future__ import annotations

import hashlib
import ipaddress
import socket
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from urllib.parse import urljoin, urlsplit, urlunsplit
from zipfile import BadZipFile, ZipFile
from collections.abc import Callable, Iterable, Sequence
from typing import Any, Protocol

import httpx
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from ebooklib import ITEM_DOCUMENT, epub
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pypdf import PdfReader

from learning_coach.ingestion import (
    MAX_EXTRACTED_CHARS,
    MAX_WEB_RESPONSE_BYTES,
    MaterialInput,
    MaterialMetadata,
    MaterialSourceType,
)
from learning_coach.media import image_bytes_content_block

MAX_ARCHIVE_MEMBERS = 2_000
MAX_ARCHIVE_UNCOMPRESSED_BYTES = 50 * 1024 * 1024
MAX_WEB_REDIRECTS = 3
MAX_IMAGE_PIXELS = 25_000_000
MAX_IMAGE_DESCRIPTION_CHARS = 8_000

CODE_LANGUAGES = {
    ".c": "c",
    ".cc": "cpp",
    ".cpp": "cpp",
    ".cs": "csharp",
    ".css": "css",
    ".go": "go",
    ".h": "c",
    ".hpp": "cpp",
    ".java": "java",
    ".js": "javascript",
    ".jsx": "javascript",
    ".json": "json",
    ".php": "php",
    ".py": "python",
    ".rb": "ruby",
    ".rs": "rust",
    ".sh": "shell",
    ".sql": "sql",
    ".toml": "toml",
    ".ts": "typescript",
    ".tsx": "typescript",
    ".xml": "xml",
    ".yaml": "yaml",
    ".yml": "yaml",
}


class MaterialLoadError(ValueError):
    """Safe user-facing ingestion error that never embeds source content."""


def _sha256(value: bytes | str) -> str:
    payload = value.encode("utf-8") if isinstance(value, str) else value
    return hashlib.sha256(payload).hexdigest()


def _source_key(material: MaterialInput) -> str:
    if material.source_url:
        parsed = urlsplit(material.source_url)
        host = (parsed.hostname or "").casefold()
        port = parsed.port
        if port and not (
            (parsed.scheme.casefold() == "http" and port == 80)
            or (parsed.scheme.casefold() == "https" and port == 443)
        ):
            host = f"{host}:{port}"
        normalized = urlunsplit(
            (
                parsed.scheme.casefold(),
                host,
                parsed.path or "/",
                parsed.query,
                "",
            )
        )
        return f"url:{normalized}"
    return f"upload:{material.source_name.casefold()}"


def _metadata(
    material: MaterialInput,
    source_type: MaterialSourceType,
    *,
    raw_content: bytes,
    location_type: str,
    location: str,
    **extra: Any,
) -> dict[str, Any]:
    source_key = _source_key(material)
    metadata = MaterialMetadata(
        source_id=_sha256(source_key),
        source_key=source_key,
        source_type=source_type,
        source_name=material.source_name,
        source_uri=material.source_uri,
        mime_type=material.mime_type or "application/octet-stream",
        content_hash=_sha256(raw_content),
        location_type=location_type,
        location=location,
        page=extra.get("page"),
        slide=extra.get("slide"),
        chapter=extra.get("chapter"),
        heading=extra.get("heading"),
        line_start=extra.get("line_start"),
        line_end=extra.get("line_end"),
    ).model_dump(exclude_none=True)
    metadata.update({key: value for key, value in extra.items() if value is not None})
    return metadata


def _decode_text(data: bytes, source_name: str) -> str:
    try:
        text = data.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise MaterialLoadError(
            f"资料 {source_name} 不是有效的 UTF-8 文本。"
        ) from exc
    return _bounded_text(text, source_name)


def _bounded_text(text: str, source_name: str) -> str:
    normalized = text.strip()
    if len(normalized) > MAX_EXTRACTED_CHARS:
        raise MaterialLoadError(
            f"资料 {source_name} 提取文本不能超过 {MAX_EXTRACTED_CHARS} 字符。"
        )
    return normalized


def _validate_archive(data: bytes, source_name: str) -> None:
    try:
        with ZipFile(BytesIO(data)) as archive:
            members = archive.infolist()
            if len(members) > MAX_ARCHIVE_MEMBERS:
                raise MaterialLoadError(f"资料 {source_name} 的压缩成员过多。")
            if sum(member.file_size for member in members) > MAX_ARCHIVE_UNCOMPRESSED_BYTES:
                raise MaterialLoadError(f"资料 {source_name} 解压后体积过大。")
    except BadZipFile as exc:
        raise MaterialLoadError(f"资料 {source_name} 不是有效的压缩文档。") from exc


def _html_documents(
    material: MaterialInput,
    raw_content: bytes,
    *,
    source_type: MaterialSourceType,
) -> list[Document]:
    html = _decode_text(raw_content, material.source_name)
    soup = BeautifulSoup(html, "html.parser")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    body = soup.body or soup
    headings = body.find_all(["h1", "h2", "h3"])
    documents: list[Document] = []
    if headings:
        for heading in headings:
            parts = [heading.get_text(" ", strip=True)]
            for sibling in heading.find_next_siblings():
                if sibling.name in {"h1", "h2", "h3"}:
                    break
                sibling_text = sibling.get_text(" ", strip=True)
                if sibling_text:
                    parts.append(sibling_text)
            text = "\n".join(part for part in parts if part).strip()
            if text:
                heading_text = parts[0]
                documents.append(
                    Document(
                        page_content=text,
                        metadata=_metadata(
                            material,
                            source_type,
                            raw_content=raw_content,
                            location_type="heading",
                            location=f"heading {heading_text}",
                            heading=heading_text,
                        ),
                    )
                )
    else:
        text = _bounded_text(body.get_text("\n", strip=True), material.source_name)
        if text:
            title = soup.title.get_text(" ", strip=True) if soup.title else "document"
            documents.append(
                Document(
                    page_content=text,
                    metadata=_metadata(
                        material,
                        source_type,
                        raw_content=raw_content,
                        location_type="document",
                        location=title,
                    ),
                )
            )
    return documents


def _default_resolver(host: str) -> list[str]:
    return sorted(
        {
            result[4][0]
            for result in socket.getaddrinfo(host, None, type=socket.SOCK_STREAM)
        }
    )


def _validate_public_url(
    url: str,
    resolver: Callable[[str], Sequence[str]],
) -> None:
    try:
        parsed = urlsplit(url)
        port = parsed.port
    except ValueError as exc:
        raise MaterialLoadError("网页 URL 无效。") from exc
    if parsed.scheme.casefold() not in {"http", "https"}:
        raise MaterialLoadError("网页资料只支持 http 或 https URL。")
    if parsed.username or parsed.password:
        raise MaterialLoadError("网页 URL 不能包含用户名或密码。")
    if not parsed.hostname:
        raise MaterialLoadError("网页 URL 缺少主机名。")
    if port is not None and not 1 <= port <= 65535:
        raise MaterialLoadError("网页 URL 端口无效。")
    try:
        direct = ipaddress.ip_address(parsed.hostname)
        addresses = [str(direct)]
    except ValueError:
        try:
            addresses = list(resolver(parsed.hostname))
        except OSError as exc:
            raise MaterialLoadError("网页主机名无法解析。") from exc
    if not addresses:
        raise MaterialLoadError("网页主机名无法解析。")
    for address in addresses:
        try:
            ip = ipaddress.ip_address(address)
        except ValueError as exc:
            raise MaterialLoadError("网页主机解析结果无效。") from exc
        if not ip.is_global:
            raise MaterialLoadError("网页资料只能访问公网地址。")


@dataclass(frozen=True)
class FetchedWebPage:
    final_url: str
    content_type: str
    content: bytes


class SafeWebFetcher:
    """Bounded HTTP fetcher that validates every redirect against SSRF."""

    def __init__(
        self,
        *,
        client: httpx.Client | None = None,
        resolver: Callable[[str], Sequence[str]] = _default_resolver,
    ) -> None:
        self._client = client
        self._resolver = resolver

    def fetch(self, url: str) -> FetchedWebPage:
        if self._client is not None:
            return self._fetch(self._client, url)
        with httpx.Client(
            follow_redirects=False,
            timeout=10.0,
            headers={"User-Agent": "learning-coach/1.0"},
        ) as client:
            return self._fetch(client, url)

    def _fetch(self, client: httpx.Client, url: str) -> FetchedWebPage:
        current = url.strip()
        for redirect_count in range(MAX_WEB_REDIRECTS + 1):
            _validate_public_url(current, self._resolver)
            try:
                with client.stream(
                    "GET", current, follow_redirects=False
                ) as response:
                    if response.status_code in {301, 302, 303, 307, 308}:
                        if redirect_count >= MAX_WEB_REDIRECTS:
                            raise MaterialLoadError("网页重定向次数过多。")
                        location = response.headers.get("location", "").strip()
                        if not location:
                            raise MaterialLoadError("网页重定向缺少目标地址。")
                        current = urljoin(current, location)
                        continue
                    try:
                        response.raise_for_status()
                    except httpx.HTTPStatusError as exc:
                        raise MaterialLoadError(
                            f"网页资料返回 HTTP {response.status_code}。"
                        ) from exc
                    declared_length = response.headers.get("content-length")
                    if declared_length and declared_length.isdigit():
                        if int(declared_length) > MAX_WEB_RESPONSE_BYTES:
                            raise MaterialLoadError("网页响应超过允许大小。")
                    content_type = (
                        response.headers.get("content-type", "")
                        .split(";", 1)[0]
                        .strip()
                        .casefold()
                    )
                    if content_type not in {
                        "application/xhtml+xml",
                        "text/html",
                        "text/markdown",
                        "text/plain",
                    }:
                        raise MaterialLoadError("网页响应不是支持的文本类型。")
                    payload = bytearray()
                    for chunk in response.iter_bytes():
                        payload.extend(chunk)
                        if len(payload) > MAX_WEB_RESPONSE_BYTES:
                            raise MaterialLoadError("网页响应超过允许大小。")
                    return FetchedWebPage(
                        final_url=str(response.url),
                        content_type=content_type,
                        content=bytes(payload),
                    )
            except httpx.HTTPError as exc:
                raise MaterialLoadError("网页资料下载失败。") from exc
        raise MaterialLoadError("网页重定向次数过多。")


class MaterialLoader(Protocol):
    """Adapter contract for one family of learning-material inputs."""

    def supports(self, material: MaterialInput) -> bool: ...

    def load(self, material: MaterialInput) -> list[Document]: ...


class MaterialLoaderRegistry:
    """Select the first explicit loader that accepts a material."""

    def __init__(self, loaders: Iterable[MaterialLoader] = ()) -> None:
        self._loaders = list(loaders)

    def register(self, loader: MaterialLoader) -> None:
        self._loaders.append(loader)

    def loader_for(self, material: MaterialInput) -> MaterialLoader:
        for loader in self._loaders:
            if loader.supports(material):
                return loader
        raise ValueError(
            f"不支持的资料格式：{material.source_name}"
        )

    def load(self, material: MaterialInput) -> list[Document]:
        loader = self.loader_for(material)
        try:
            documents = loader.load(material)
        except MaterialLoadError:
            raise
        except Exception as exc:
            raise MaterialLoadError(
                f"无法解析资料 {material.source_name}。"
            ) from exc
        if not documents:
            raise ValueError(f"资料没有可提取的文本：{material.source_name}")
        return documents


class PdfMaterialLoader:
    def supports(self, material: MaterialInput) -> bool:
        return material.suffix == ".pdf" or material.mime_type == "application/pdf"

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        reader = PdfReader(BytesIO(material.data))
        if reader.is_encrypted:
            raise MaterialLoadError(f"资料 {material.source_name} 受密码保护。")
        documents: list[Document] = []
        total = 0
        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if not text:
                continue
            total += len(text)
            if total > MAX_EXTRACTED_CHARS:
                raise MaterialLoadError(
                    f"资料 {material.source_name} 提取文本不能超过 {MAX_EXTRACTED_CHARS} 字符。"
                )
            documents.append(
                Document(
                    page_content=text,
                    metadata=_metadata(
                        material,
                        "pdf",
                        raw_content=material.data,
                        location_type="page",
                        location=f"page {page_number}",
                        page=page_number,
                    ),
                )
            )
        return documents


class DocxMaterialLoader:
    _MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"

    def supports(self, material: MaterialInput) -> bool:
        return material.suffix == ".docx" or material.mime_type == self._MIME

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        _validate_archive(material.data, material.source_name)
        document = DocxDocument(BytesIO(material.data))
        documents: list[Document] = []
        total = 0
        current_heading: str | None = None
        for paragraph_number, paragraph in enumerate(document.paragraphs, start=1):
            text = paragraph.text.strip()
            if not text:
                continue
            total += len(text)
            if total > MAX_EXTRACTED_CHARS:
                raise MaterialLoadError(
                    f"资料 {material.source_name} 提取文本不能超过 {MAX_EXTRACTED_CHARS} 字符。"
                )
            style_name = paragraph.style.name if paragraph.style is not None else ""
            if style_name.casefold().startswith("heading"):
                current_heading = text
            documents.append(
                Document(
                    page_content=text,
                    metadata=_metadata(
                        material,
                        "docx",
                        raw_content=material.data,
                        location_type="paragraph",
                        location=f"paragraph {paragraph_number}",
                        paragraph=paragraph_number,
                        heading=current_heading,
                    ),
                )
            )
        return documents


class PptxMaterialLoader:
    _MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def supports(self, material: MaterialInput) -> bool:
        return material.suffix == ".pptx" or material.mime_type == self._MIME

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        _validate_archive(material.data, material.source_name)
        presentation = Presentation(BytesIO(material.data))
        documents: list[Document] = []
        total = 0
        for slide_number, slide in enumerate(presentation.slides, start=1):
            parts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            text = "\n".join(parts)
            if not text:
                continue
            total += len(text)
            if total > MAX_EXTRACTED_CHARS:
                raise MaterialLoadError(
                    f"资料 {material.source_name} 提取文本不能超过 {MAX_EXTRACTED_CHARS} 字符。"
                )
            documents.append(
                Document(
                    page_content=text,
                    metadata=_metadata(
                        material,
                        "pptx",
                        raw_content=material.data,
                        location_type="slide",
                        location=f"slide {slide_number}",
                        slide=slide_number,
                    ),
                )
            )
        return documents


class EpubMaterialLoader:
    def supports(self, material: MaterialInput) -> bool:
        return material.suffix == ".epub" or material.mime_type == "application/epub+zip"

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        _validate_archive(material.data, material.source_name)
        book = epub.read_epub(BytesIO(material.data))
        documents: list[Document] = []
        total = 0
        for item in book.get_items_of_type(ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            for element in soup(["script", "style", "noscript"]):
                element.decompose()
            text = soup.get_text("\n", strip=True)
            if not text:
                continue
            total += len(text)
            if total > MAX_EXTRACTED_CHARS:
                raise MaterialLoadError(
                    f"资料 {material.source_name} 提取文本不能超过 {MAX_EXTRACTED_CHARS} 字符。"
                )
            heading_node = soup.find(["h1", "h2", "h3"])
            chapter = (
                getattr(item, "title", None)
                or (heading_node.get_text(" ", strip=True) if heading_node else None)
                or item.get_name()
            )
            documents.append(
                Document(
                    page_content=text,
                    metadata=_metadata(
                        material,
                        "epub",
                        raw_content=material.data,
                        location_type="chapter",
                        location=f"chapter {chapter}",
                        chapter=str(chapter),
                    ),
                )
            )
        return documents


class HtmlMaterialLoader:
    def supports(self, material: MaterialInput) -> bool:
        return material.suffix in {".htm", ".html"} or material.mime_type == "text/html"

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        return _html_documents(material, material.data, source_type="html")


class WebMaterialLoader:
    def __init__(self, fetcher: SafeWebFetcher | None = None) -> None:
        self._fetcher = fetcher or SafeWebFetcher()

    def supports(self, material: MaterialInput) -> bool:
        return material.source_url is not None

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.source_url is not None
        fetched = self._fetcher.fetch(material.source_url)
        fetched_material = MaterialInput(
            source_name=material.source_name,
            mime_type=fetched.content_type,
            source_url=fetched.final_url,
        )
        if fetched.content_type in {"text/plain", "text/markdown"}:
            text = _decode_text(fetched.content, material.source_name)
            return [
                Document(
                    page_content=text,
                    metadata=_metadata(
                        fetched_material,
                        "web",
                        raw_content=fetched.content,
                        location_type="url",
                        location=fetched.final_url,
                    ),
                )
            ]
        return _html_documents(
            fetched_material,
            fetched.content,
            source_type="web",
        )


def _message_text(message: Any) -> str:
    text = getattr(message, "text", None)
    if isinstance(text, str) and text.strip():
        return text.strip()
    content = getattr(message, "content", message)
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for block in content:
            if isinstance(block, str):
                parts.append(block)
            elif isinstance(block, dict) and isinstance(block.get("text"), str):
                parts.append(block["text"])
        return "".join(parts).strip()
    return ""


class ImageMaterialLoader:
    _SUFFIXES = {".gif", ".jpeg", ".jpg", ".png", ".webp"}

    def __init__(self, *, model: Any | None, accepts_images: bool) -> None:
        self._model = model
        self._accepts_images = accepts_images

    def supports(self, material: MaterialInput) -> bool:
        return material.suffix in self._SUFFIXES or material.mime_type.startswith("image/")

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        if not self._accepts_images or self._model is None:
            raise MaterialLoadError("当前主模型不支持图片资料摄取。")
        try:
            with Image.open(BytesIO(material.data)) as image:
                width, height = image.size
                image.verify()
        except (OSError, UnidentifiedImageError) as exc:
            raise MaterialLoadError(
                f"无法解析图片资料 {material.source_name}。"
            ) from exc
        if width * height > MAX_IMAGE_PIXELS:
            raise MaterialLoadError(f"图片资料 {material.source_name} 像素过大。")
        image_block = image_bytes_content_block(material.data, material.mime_type)
        prompt = (
            "请把这张学习资料图片转换为可检索文本。先逐字记录清晰可见的文字，"
            "再描述图表、流程、公式、代码和关键关系。不要补充图片中不存在的信息，"
            f"总长度不超过 {MAX_IMAGE_DESCRIPTION_CHARS} 个字符。"
        )
        try:
            response = self._model.invoke(
                [
                    HumanMessage(
                        content=[
                            {"type": "text", "text": prompt},
                            image_block,
                        ]
                    )
                ]
            )
        except Exception as exc:
            raise MaterialLoadError(
                f"图片资料 {material.source_name} 解析失败。"
            ) from exc
        description = _message_text(response)
        if not description:
            raise MaterialLoadError(f"图片资料 {material.source_name} 没有可提取内容。")
        description = description[:MAX_IMAGE_DESCRIPTION_CHARS].strip()
        return [
            Document(
                page_content=description,
                metadata=_metadata(
                    material,
                    "image",
                    raw_content=material.data,
                    location_type="image",
                    location=material.source_name,
                    width=width,
                    height=height,
                ),
            )
        ]


class CodeMaterialLoader:
    def supports(self, material: MaterialInput) -> bool:
        return material.suffix in CODE_LANGUAGES

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        text = _decode_text(material.data, material.source_name)
        if not text:
            return []
        lines = text.splitlines()
        line_end = max(1, len(lines))
        language = CODE_LANGUAGES[material.suffix]
        return [
            Document(
                page_content=text,
                metadata=_metadata(
                    material,
                    "code",
                    raw_content=material.data,
                    location_type="lines",
                    location=f"lines 1-{line_end}",
                    line_start=1,
                    line_end=line_end,
                    language=language,
                ),
            )
        ]


class TextMaterialLoader:
    _SUFFIXES = {".md", ".markdown", ".txt"}
    _MIME_TYPES = {"text/markdown", "text/plain"}

    def supports(self, material: MaterialInput) -> bool:
        return material.suffix in self._SUFFIXES or material.mime_type in self._MIME_TYPES

    def load(self, material: MaterialInput) -> list[Document]:
        assert material.data is not None
        text = _decode_text(material.data, material.source_name)
        if not text:
            return []
        return [
            Document(
                page_content=text,
                metadata=_metadata(
                    material,
                    "text",
                    raw_content=material.data,
                    location_type="document",
                    location="document",
                ),
            )
        ]


def default_loader_registry(
    *,
    web_fetcher: SafeWebFetcher | None = None,
    image_model: Any | None = None,
    accepts_images: bool = False,
) -> MaterialLoaderRegistry:
    """Build the bounded multi-format loader registry."""

    return MaterialLoaderRegistry(
        [
            WebMaterialLoader(web_fetcher),
            PdfMaterialLoader(),
            DocxMaterialLoader(),
            PptxMaterialLoader(),
            EpubMaterialLoader(),
            HtmlMaterialLoader(),
            ImageMaterialLoader(
                model=image_model,
                accepts_images=accepts_images,
            ),
            CodeMaterialLoader(),
            TextMaterialLoader(),
        ]
    )
