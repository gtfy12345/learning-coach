from io import BytesIO

import httpx
import pytest
from docx import Document as DocxDocument
from ebooklib import epub
from langchain_core.messages import AIMessage
from PIL import Image
from pptx import Presentation

from learning_coach.ingestion import MAX_WEB_RESPONSE_BYTES, MaterialInput
from learning_coach.loaders import (
    ImageMaterialLoader,
    SafeWebFetcher,
    WebMaterialLoader,
    default_loader_registry,
)


def _minimal_pdf(text: str) -> bytes:
    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    stream = f"BT /F1 12 Tf 72 720 Td ({escaped}) Tj ET".encode("ascii")
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length "
        + str(len(stream)).encode("ascii")
        + b" >>\nstream\n"
        + stream
        + b"\nendstream",
    ]
    payload = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, body in enumerate(objects, start=1):
        offsets.append(len(payload))
        payload.extend(f"{number} 0 obj\n".encode("ascii"))
        payload.extend(body)
        payload.extend(b"\nendobj\n")
    xref = len(payload)
    payload.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    payload.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        payload.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    payload.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref}\n%%EOF\n"
        ).encode("ascii")
    )
    return bytes(payload)


def _docx_bytes() -> bytes:
    document = DocxDocument()
    document.add_heading("Reducer 课程", level=1)
    document.add_paragraph("Reducer 合并并行状态更新。")
    output = BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "LangGraph 课程"
    slide.placeholders[1].text = "条件边根据结构化状态选择下一节点。"
    output = BytesIO()
    presentation.save(output)
    return output.getvalue()


def _epub_bytes() -> bytes:
    book = epub.EpubBook()
    book.set_identifier("learning-coach")
    book.set_title("学习教练")
    book.set_language("zh")
    chapter = epub.EpubHtml(title="第一章", file_name="chapter.xhtml", lang="zh")
    chapter.content = "<h1>第一章</h1><p>Runnable 可以组合成教学任务。</p>"
    book.add_item(chapter)
    book.spine = ["nav", chapter]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    output = BytesIO()
    epub.write_epub(output, book)
    return output.getvalue()


@pytest.mark.parametrize(
    ("material", "expected_text", "location_key", "location_value"),
    [
        (
            MaterialInput("paper.pdf", "application/pdf", data=_minimal_pdf("Reducer paper")),
            "Reducer paper",
            "page",
            1,
        ),
        (
            MaterialInput(
                "course.docx",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                data=_docx_bytes(),
            ),
            "Reducer 合并",
            "paragraph",
            2,
        ),
        (
            MaterialInput(
                "slides.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                data=_pptx_bytes(),
            ),
            "条件边",
            "slide",
            1,
        ),
        (
            MaterialInput("book.epub", "application/epub+zip", data=_epub_bytes()),
            "Runnable",
            "chapter",
            "第一章",
        ),
    ],
)
def test_document_loaders_preserve_format_location(
    material: MaterialInput,
    expected_text: str,
    location_key: str,
    location_value: int | str,
) -> None:
    documents = default_loader_registry().load(material)

    matching = [document for document in documents if expected_text in document.page_content]
    assert matching
    assert matching[0].metadata[location_key] == location_value
    assert matching[0].metadata["source_name"] == material.source_name
    assert len(matching[0].metadata["source_id"]) == 64
    assert len(matching[0].metadata["content_hash"]) == 64


def test_html_text_markdown_and_code_loaders_keep_structure() -> None:
    registry = default_loader_registry()
    html = registry.load(
        MaterialInput(
            "lesson.html",
            "text/html",
            data=(
                b"<html><head><title>LCEL Course</title><script>secret()</script></head>"
                b"<body><h1>Runnable</h1><p>Sequence composes steps.</p></body></html>"
            ),
        )
    )
    markdown = registry.load(
        MaterialInput("notes.md", "text/markdown", data="# Graph\nReducer 合并状态".encode())
    )
    text = registry.load(
        MaterialInput("notes.txt", "text/plain", data="学习资料".encode())
    )
    code = registry.load(
        MaterialInput(
            "graph.py",
            "text/x-python",
            data=b"def route(score):\n    return 'finish' if score >= 80 else 'retry'\n",
        )
    )

    assert html[0].metadata["heading"] == "Runnable"
    assert "secret" not in " ".join(document.page_content for document in html)
    assert markdown[0].metadata["location_type"] == "document"
    assert text[0].page_content == "学习资料"
    assert code[0].metadata["line_start"] == 1
    assert code[0].metadata["line_end"] == 2
    assert code[0].metadata["language"] == "python"


def test_corrupt_or_empty_document_fails_without_exposing_content() -> None:
    registry = default_loader_registry()

    with pytest.raises(ValueError, match="无法解析资料 paper.pdf") as exc_info:
        registry.load(
            MaterialInput("paper.pdf", "application/pdf", data=b"not a pdf secret-body")
        )

    assert "secret-body" not in str(exc_info.value)


def test_web_loader_fetches_public_html_and_preserves_final_url() -> None:
    def handle(request: httpx.Request) -> httpx.Response:
        assert request.url == httpx.URL("https://example.com/course")
        return httpx.Response(
            200,
            headers={"content-type": "text/html; charset=utf-8"},
            content=(
                "<html><body><h1>课程网页</h1><p>Reducer 合并状态。</p>"
                "<script>do_not_index()</script></body></html>"
            ).encode(),
            request=request,
        )

    client = httpx.Client(transport=httpx.MockTransport(handle), follow_redirects=False)
    fetcher = SafeWebFetcher(
        client=client,
        resolver=lambda host: ["93.184.216.34"],
    )
    loader = WebMaterialLoader(fetcher)

    documents = loader.load(
        MaterialInput(
            "课程网页",
            "text/html",
            source_url="https://example.com/course",
        )
    )

    assert "Reducer" in documents[0].page_content
    assert "do_not_index" not in documents[0].page_content
    assert documents[0].metadata["source_type"] == "web"
    assert documents[0].metadata["source_uri"] == "https://example.com/course"
    assert documents[0].metadata["heading"] == "课程网页"
    client.close()


@pytest.mark.parametrize("address", ["127.0.0.1", "10.0.0.8", "169.254.169.254", "::1"])
def test_web_fetcher_blocks_non_public_addresses_before_request(address: str) -> None:
    called = False

    def handle(request: httpx.Request) -> httpx.Response:
        nonlocal called
        called = True
        return httpx.Response(200, content=b"never", request=request)

    client = httpx.Client(transport=httpx.MockTransport(handle), follow_redirects=False)
    fetcher = SafeWebFetcher(client=client, resolver=lambda host: [address])

    with pytest.raises(ValueError, match="公网地址"):
        fetcher.fetch("https://example.com/private")
    assert called is False
    client.close()


def test_web_fetcher_revalidates_redirect_and_bounds_response() -> None:
    def redirect_to_private(request: httpx.Request) -> httpx.Response:
        return httpx.Response(
            302,
            headers={"location": "http://127.0.0.1/admin"},
            request=request,
        )

    redirect_client = httpx.Client(
        transport=httpx.MockTransport(redirect_to_private), follow_redirects=False
    )
    fetcher = SafeWebFetcher(
        client=redirect_client,
        resolver=lambda host: ["93.184.216.34"],
    )
    with pytest.raises(ValueError, match="公网地址"):
        fetcher.fetch("https://example.com/redirect")
    redirect_client.close()

    def oversized(request: httpx.Request) -> httpx.Response:
        return httpx.Response(
            200,
            headers={"content-type": "text/html"},
            content=b"x" * (MAX_WEB_RESPONSE_BYTES + 1),
            request=request,
        )

    oversized_client = httpx.Client(
        transport=httpx.MockTransport(oversized), follow_redirects=False
    )
    oversized_fetcher = SafeWebFetcher(
        client=oversized_client,
        resolver=lambda host: ["93.184.216.34"],
    )
    with pytest.raises(ValueError, match="网页响应"):
        oversized_fetcher.fetch("https://example.com/huge")
    oversized_client.close()


def test_web_fetcher_stops_streaming_after_response_limit() -> None:
    class CountingStream(httpx.SyncByteStream):
        def __init__(self) -> None:
            self.chunks_read = 0

        def __iter__(self):
            for _ in range(4):
                self.chunks_read += 1
                yield b"x" * (MAX_WEB_RESPONSE_BYTES // 2)

    stream = CountingStream()

    def oversized_stream(request: httpx.Request) -> httpx.Response:
        return httpx.Response(
            200,
            headers={"content-type": "text/html"},
            stream=stream,
            request=request,
        )

    client = httpx.Client(
        transport=httpx.MockTransport(oversized_stream), follow_redirects=True
    )
    fetcher = SafeWebFetcher(
        client=client,
        resolver=lambda host: ["93.184.216.34"],
    )

    with pytest.raises(ValueError, match="网页响应"):
        fetcher.fetch("https://example.com/stream")

    assert stream.chunks_read == 3
    client.close()


class _VisionModel:
    def __init__(self) -> None:
        self.calls: list[object] = []

    def invoke(self, messages: object) -> AIMessage:
        self.calls.append(messages)
        return AIMessage(content="图片文字：StateGraph。图中展示诊断到补救的有向流程。")


def _png_bytes() -> bytes:
    output = BytesIO()
    Image.new("RGB", (16, 12), color=(255, 255, 255)).save(output, format="PNG")
    return output.getvalue()


def test_image_loader_uses_one_vision_call_and_records_dimensions() -> None:
    model = _VisionModel()
    loader = ImageMaterialLoader(model=model, accepts_images=True)
    material = MaterialInput("workflow.png", "image/png", data=_png_bytes())

    documents = loader.load(material)

    assert len(model.calls) == 1
    message = model.calls[0][0]
    assert message.content[1]["type"] == "image"
    assert "StateGraph" in documents[0].page_content
    assert documents[0].metadata["source_type"] == "image"
    assert documents[0].metadata["width"] == 16
    assert documents[0].metadata["height"] == 12


def test_image_loader_rejects_missing_capability_and_corrupt_image() -> None:
    material = MaterialInput("workflow.png", "image/png", data=_png_bytes())
    with pytest.raises(ValueError, match="不支持图片资料"):
        ImageMaterialLoader(model=_VisionModel(), accepts_images=False).load(material)

    with pytest.raises(ValueError, match="无法解析图片资料"):
        ImageMaterialLoader(model=_VisionModel(), accepts_images=True).load(
            MaterialInput("broken.png", "image/png", data=b"not-an-image")
        )
